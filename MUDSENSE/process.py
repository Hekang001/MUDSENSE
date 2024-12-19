import pandas as pd
import numpy as np
# import cv2
# import pytesseract
# import re
import requests
import eml_parser
import re
import bs4
import base64
import datetime
import json
from regipy import RegistryHive
from tqdm import tqdm
import os
import zipfile
import shutil
import urllib
from pptx import Presentation
import base64
import requests
import json
from docx import Document
from os.path import basename
import re
import os
import time
import jsonlines

# xyy
access_token = '24.1a0a4622c1ece1e9cd24b1379f9586fe.2592000.1700296541.282335-38817714'
# # hk
# access_token = '24.f3598648a9da3989fbc6306bc821003f.2592000.1701326008.282335-42095745'

def get_access_token():
    """
    使用 AK，SK 生成鉴权签名（Access Token）
    :return: access_token，或是None(如果错误)
    """
    # # xyy
    # API_KEY = "cT3xO7tLD2KDFSwYb1ZesTCA"
    # SECRET_KEY = "q2EtxPIGVOBwAdHOHjSP90UY8jtydM94"
    # hk
    API_KEY = "OfTw9e9pmr3fZFGVbiZsEGeG"
    SECRET_KEY = "v3rx87R4RvK0GTTw1Mth8KmDUCTVRxz3"
    url = "https://aip.baidubce.com/oauth/2.0/token"
    params = {"grant_type": "client_credentials", "client_id": API_KEY, "client_secret": SECRET_KEY}
    return str(requests.post(url, params=params).json().get("access_token"))


'''预测.xlsx、.et形式表格'''
def predict_form(path):
    result = []
    count = 0
    while True:
        try:
            df = pd.read_excel(path, sheet_name=count)
            data = df.values
            keys = df.columns
            for i in range(data.shape[0]):
                a_dict = {keys[j]: data[i][j] for j in range(data.shape[1])}
                result.append(a_dict)
        except:
            break
        count += 1
    return result


'''填补空缺点位'''
def process_pos(bitwise_and):
    # 点横纵坐标
    all = np.where(bitwise_and == 255)
    heng_point = list(sorted(set(all[0])))
    zong_point = list(sorted(set(all[1])))

    for i in range(len(heng_point)):
        for j in range(len(zong_point)):
            bitwise_and[heng_point[i]][zong_point[j]] = 255
    return bitwise_and, heng_point, zong_point


'''预测图片形式表格'''
def predict_picture_form(path):
    result = []

    # return result
    request_url = "https://aip.baidubce.com/rest/2.0/ocr/v1/table"
    # 二进制方式打开图片文件
    f = open(path, 'rb')
    img = base64.b64encode(f.read())

    params = {"image": img}
    request_url = request_url + "?access_token=" + access_token
    headers = {'content-type': 'application/x-www-form-urlencoded'}

    response = requests.post(request_url, data=params, headers=headers)
    response_json = response.json()
    if response and response_json['table_num']:
        table = response_json['tables_result'][0]['body']
        max_col = table[0]['col_end']
        begin_col = table[0]['col_end']
        for i in range(1, len(table)):
            if table[i]['col_end'] == begin_col:
                break
            if table[i]['col_end'] > max_col:
                max_col = table[i]['col_end']
        max_row = int(len(table) / max_col)

        keys = []
        for i in range(max_row):
            if i != 0:
                a_dict = {}
            for j in range(max_col):
                if i == 0:
                    keys.append(table[i * max_col + j]['words'])
                else:
                    a_dict[keys[j]] = table[i * max_col + j]['words']
            if i != 0:
                result.append(a_dict)
    return result


'''预测邮件中的表格'''
def predict_mail_form(path):
    with open(path, 'rb') as f:
        a = f.read()
    # decode_email_b参数
    # para1:eml文件二进制读出来的内容
    # para2:表示邮件解析返回的数据结构里是否包含邮件的原始内容，也就是html内容，我这里选True，这样内容可以直接放到网页中显示。
    # 解析返回的数据结构中是否包含附件内容。
    # eml_parser==1.10.0
    eml = eml_parser.eml_parser.decode_email_b(a, True, True)
    # eml = eml_parser.decode_email_b(a, True, True)


    title = eml["header"]["subject"]
    reciver = ','.join(eml["header"]["to"])
    sender = eml["header"]["from"]
    date = eml["header"]['date']
    mail_info = [{"title": title, "reciver": reciver, "sender": sender, "date": date}]

    for i in eml["body"]:
        # text/html是带html标签，可以直接放到网页里面展示
        # 还有类型text/plain,正文部分的文字信息
        if i["content_type"] == "text/html":
            # 下两行是对html样式进行了调整以适配前端显示，可自行调整
            body = i["content"].replace('\n', '').replace('position:absolute;', '')
            body = re.sub('body\s*{.*?}', '', body)
            break

    # .eml中直接出现的表格
    table = bs4.BeautifulSoup(body, features="lxml").find("table")
    table_data_list = [[cell.text.strip() for cell in row.find_all("td")] for row in table.find_all("tr")]
    table_data = []
    for i in range(1, len(table_data_list)):
        table_data.append({table_data_list[0][j]: table_data_list[i][j] for j in range(len(table_data_list[0]))})

    # 除表格外内容
    email_text=""
    for i in eml["body"]:
        try:
            if i['content_type'] == 'text/plain':
                email_text = i['content']
                break
            elif i['content_type'] == 'text/html':
                soup = bs4.BeautifulSoup(i['content'], 'lxml')
                p_data = soup.find_all('p')
                ree = []
                for i in p_data:
                    res = i.get_text()
                    ree.append(res)
                email_text = '\n'.join(ree)
        except Exception as e:
            email_text = i["content"].replace('\r\n', '')

    # .eml中附件表格
    out_path = "temp.xlsx"
    if eml.get('attachment'):
        for i in eml['attachment']:
            x = base64.b64decode(i['raw'])
            with open(out_path, 'wb') as f:
                f.write(x)
        excel_result = predict_form(out_path)
    else:
        excel_result = []
    os.remove(out_path)
    return table_data + excel_result + mail_info,email_text


'''php中Key是否所需'''
def php_valid(string):
    valid_key = ["NAME", "USER", "PASSWORD", "HOST", "KEY", "SALT", "AUTH"]
    for key in valid_key:
        if key in string:
            return True
    return False


'''识别path路径图片中的文字,need_list:结果list是否需要拼接'''
def get_img_text(path,need_list):
    '''通用文字识别（高精度版）'''
    request_url = "https://aip.baidubce.com/rest/2.0/ocr/v1/accurate_basic"
    # 二进制方式打开图片文件
    f = open(path, 'rb')
    img = base64.b64encode(f.read())

    params = {"image": img}
    request_url = request_url + "?access_token=" + access_token
    headers = {'content-type': 'application/x-www-form-urlencoded'}
    response = requests.post(request_url, data=params, headers=headers)
    result_list=[]
    result_str=""
    if response:
        try:
            result_list = response.json()['words_result']
            if not need_list:
                for result in result_list:
                    result_str += result['words']
                    result_str += "\n"
        except:
            print("error!!!!!!!!!!")
    if need_list:
        return result_list
    else:
        return result_str


'''ocr、预测图片中的php代码'''
def predict_picture_php_ocr(path):
    # ocr
    result_list=get_img_text(path,need_list=True)
    # php
    if result_list[0]["words"] == '<?php':
        result = {}
        i = 1
        while i < len(result_list):
            sent = result_list[i]['words']
            if sent.startswith("define"):
                sent_list = sent.split("'")
                if php_valid(sent_list[1]):
                    if sent_list[-1].startswith(");") or sent_list[-1].startswith(" );"):
                        result[sent_list[1]] = sent_list[-2]
                    else:
                        i += 1
                        next_sent_list = result_list[i]['words'].strip().split("'")
                        result[sent_list[1]] = next_sent_list[1]
            i += 1
        return True,[result]
    # 其他文本
    else:
        result_str=""
        for result in result_list:
            result_str += result['words']
            result_str += "\n"
        return False,result_str


'''识别pdf中文字'''
def get_pdf_text(path):
    '''通用文字识别（高精度版）'''
    request_url = "https://aip.baidubce.com/rest/2.0/ocr/v1/accurate_basic"
    # 二进制方式打开图片文件
    f = open(path, 'rb')
    pdf_content = base64.b64encode(f.read()).decode("utf8")
    pdf_content=urllib.parse.quote_plus(pdf_content)
    request_url = request_url + "?access_token=" + access_token
    headers = {
        'Content-Type': 'application/x-www-form-urlencoded',
        'Accept': 'application/json'
    }

    response = requests.request("POST", request_url, headers=headers, data="pdf_file="+pdf_content)
    result_str = ""
    if response:
        try:
            result_list = response.json()['words_result']
            for result in result_list:
                result_str += result['words']
                result_str += "\n"
        except:
            print("error!!!!!!!!!!")
    return result_str


'''分离mysql命令中的-x xxx, 得到对应值xxx'''
def split_value(command, string):
    if command in string:
        temp_ip = string.split(command)
        temp_ip = temp_ip[1].strip()
        if " " in temp_ip:
            value = temp_ip.split(" ")[0]
        else:
            value = temp_ip
        return value
    return ""


'''抽取.bash_history'''
def preict_command(path):
    with open(path, 'r') as f:
        datas = f.read()
    data_list = datas.strip().split("\n")

    ip_pattern = r"\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}"
    port_pattern_1 = r"\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}:\d{1,5}"
    port_pattern_2 = r"-p \d{1,5}"

    result = []
    for data in data_list:
        a_dict = {}
        # mysql特殊处理
        if data.startswith("mysql"):
            ip = split_value("-h", data)
            if ip != "":
                a_dict["ip"] = ip
            username = split_value("-u", data)
            if username != "":
                a_dict["username"] = username
            password = split_value("-p", data)
            if password != "":
                a_dict["password"] = password
        else:
            # ip
            ip_result = re.findall(ip_pattern, data)
            if ip_result:
                a_dict["ip"] = ip_result[0]
            # port, 格式ip:port
            port_result = re.findall(port_pattern_1, data)
            if port_result:
                a_dict["port"] = port_result[0].split(":")[1]
            else:
                # port, 格式xx -p port
                port_result = re.findall(port_pattern_2, data)
                if port_result:
                    a_dict["port"] = port_result[0].split(" ")[1]

        if a_dict != {}:
            result.append(a_dict)
    return result


'''处理txt文本'''
def predict_txt(path):
    ip_pattern = r"\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}"
    phone_pattern = r"^1[34578]\d{9}$"
    with open(path, "r", encoding="utf-8") as file:
        content = file.read()
    cont_list=content.split("\n\n")
    cont_list=[cont for cont in cont_list if cont]
    # 提取
    result = []
    info_type = ["ip", "mail", "account", "password","phone_number"]
    other_text=""
    for cont in cont_list:
        info_list=cont.split("\n")
        info_list = [info for info in info_list if info]
        cont_result={"ip":[],"account":[],"password":[],"mail":[],"phone_number":[]}
        flag=0
        for info in info_list:
            if re.match(ip_pattern,info):
                cont_result["ip"].append(info)
                # 前一行为ip
                flag=1
            elif info=="root" or flag==1:
                cont_result["account"].append(info)
                # 前一行为账号
                flag=2
            elif flag==2:
                cont_result["password"].append(info)
            elif ".com" in info:
                cont_result["mail"].append(info)
            elif re.match(phone_pattern,info):
                cont_result["phone_number"].append(info)
            else:
                other_text+=info
                other_text+="\n"
        # 处理
        # 哪一项包含多个
        flag=0
        if len(cont_result["ip"])>1: flag=1
        elif len(cont_result["mail"])>1: flag=2
        elif len(cont_result["account"])>1: flag=3
        elif len(cont_result["password"])>1: flag=4
        elif len(cont_result["password"]) > 1:flag = 5
        # 重复词典，替换类型相同信息
        an_dict={key:cont_result[key][0] for key in cont_result.keys() if cont_result[key]!=[]}
        if flag==0:
            an_result=[an_dict]
        else:
            an_result=[an_dict.copy() for i in range(len(cont_result[info_type[flag-1]]))]
            for j in range(len(cont_result[info_type[flag-1]])-1):
                an_result[j+1][info_type[flag-1]]=cont_result[info_type[flag-1]][j+1]
        result+=an_result
    return result,other_text


''' .hiv文件解析函数 .hiv 文件是 Windows 注册表中的一个二进制文件，用于存储特定注册表项及其关联的数据。'''
def parse_and_save_hiv(hiv_file_path):
    abspath=os.path.abspath(hiv_file_path)
    hive = RegistryHive(abspath)
    # 可以创建一个名为hive的注册表对象，并加载指定路径hiv_file_path的注册表文件
    output_list,hiv_dict_list,hiv_dict = [],[],{}
    # 遍历注册表键并添加到输出列表中
    for key in tqdm(hive.recurse_subkeys()):
        output_list.append("Key Path: " + key.path)
        output_list.append("Number of Values: " + str(len(key.values)))
        # 写入字典
        hiv_dict['Key Path'] = str(key.path)
        hiv_dict['Number of Values'] = str(len(key.values))
        # 遍历注册表值并添加到输出列表中
        for value in key.values:
            value_name = value.name
            value_type = str(value.value_type)
            value_data = str(value.value)
            output_list.append("Value Name: " + value_name)
            output_list.append("Value Type: " + value_type)
            output_list.append("Value Data: " + value_data)
            output_list.append("---")
            # 写入字典
            hiv_dict['Value Name'] = value_name
            hiv_dict['Value Type'] = value_type
            hiv_dict['Value Data'] = value_data
            # 将字典导入列表，并新建一个字典
            hiv_dict_list.append(hiv_dict)
            hiv_dict = {}
    return output_list,hiv_dict_list


'''抽取pptx/docx中的表格'''
def extract_pptx_docx(tables):
    result = []
    for table in tables:
        keys = []
        for i in range(len(table.rows)):
            row = table.rows[i]
            a_result = {}
            for j in range(1, len(row.cells)):
                # 键值信息
                cell_text = row.cells[j].text
                if i == 0:
                    keys.append(cell_text)
                else:
                    a_result[keys[j - 1]] = cell_text
            if a_result != {}:
                result.append(a_result)
    return result


'''pptx表格抽取，识别文字(含图片)'''
def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    ppt_text=""
    tables=[]
    for slide in presentation.slides:
        for shape in slide.shapes:
            # 文本
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = ' '.join(run.text for run in paragraph.runs)
                    ppt_text+=text
                    ppt_text+="\n"
            # 表格
            if shape.has_table:
                tables.append(shape.table)
            # 图片
            try:
                if "image" in shape.image.content_type:
                    # 临时保存并识别
                    imgName = shape.image.filename
                    with open(imgName, "wb") as f:
                        f.write(shape.image.blob)
                    time.sleep(1)
                    img_text=get_img_text(imgName,need_list=False)
                    ppt_text+=img_text
                    os.remove(imgName)
            except:
                continue
    result=extract_pptx_docx(tables)
    return result,ppt_text


'''docx表格抽取，识别文字(含图片)'''
def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    result = list()
    pattern = re.compile('rId\d+')
    for graph in doc.paragraphs:
        para_result = list()
        for run in graph.runs:
            if run.text != '':
                para_result.append(run.text.strip())
            else:
                contentID = pattern.search(run.element.xml).group(0)
                try:
                    contentType = doc.part.related_parts[contentID].content_type
                except KeyError as e:
                    continue
                if not contentType.startswith('image'):
                    continue
                imgName = basename(doc.part.related_parts[contentID].partname)
                imgData = doc.part.related_parts[contentID].blob
                # 临时保存并识别
                with open(imgName, "wb") as f:
                    f.write(imgData)
                time.sleep(1)
                img_text = get_img_text(imgName,need_list=False)
                para_result.append(img_text.strip())
                os.remove(imgName)
        if para_result != []:
            result.append("".join(para_result))
    docx_text = "\n".join(result)
    # 抽取表格
    result=extract_pptx_docx(doc.tables)
    return result,docx_text


def predict_test(input_path, extract_path, text_path):
    # 获取所有文件路径
    file_ls = []
    for root, dirs, files in os.walk(input_path):
        root_file_ls = [os.path.join(root, file) for file in files]
        file_ls += root_file_ls

    result = []
    text_result = []
    # .pptx文本处理
    for file in file_ls:
        # .docx表格抽取，文本处理
        if file.endswith("麒麟SSL+VPN+Windows客户端使用手册.docx"):
            docx_result, docx_text = extract_text_from_docx(file)
            text_result.append({file: docx_text})
            if docx_result != []:
                result.append({file: docx_result})


def predict(input_path,extract_path,text_path):
    # 获取所有文件路径
    file_ls = []
    for root, dirs, files in os.walk(input_path):
        root_file_ls = [os.path.join(root, file) for file in files]
        file_ls += root_file_ls

    result = []
    text_result=[]
    for file in file_ls:
        # 是否处理
        flag=1
        # .xlsx/.et表格抽取
        if file[-5:] == ".xlsx" or file[-3:] == ".et":
            a_result = predict_form(file)
            if a_result != []:
                result.append({file: a_result})
        # .png/.jpg: 表格抽取、php抽取、其他文本处理
        elif file[-4:] == ".png" or file[-4:] == ".jpg":
            # .png/.jpg表格抽取
            a_result = predict_picture_form(file)
            if a_result != []:
                result.append({file: a_result})
            else:
                # php抽取、其他文本处理
                is_php, a_result = predict_picture_php_ocr(file)
                if is_php and a_result != []:
                    result.append({file: a_result})
                elif not is_php and a_result!="":
                    text_result.append({file:a_result})
        # .eml中表格抽取、其他文本处理
        elif file[-4:] == ".eml":
            a_result,email_text = predict_mail_form(file)
            if a_result != []:
                result.append({file: a_result})
            if email_text!="":
                text_result.append({file:email_text})
        # txt规则抽取，其他文本处理
        elif file[-4:]==".txt":
            a_result,other_text = predict_txt(file)
            if a_result != []:
                result.append({file: a_result})
            if other_text!="":
                text_result.append({file:other_text})
        # linux文件夹中文件抽取
        elif file.split("/")[2] == "linux":
            with open(file, 'r', encoding="utf-8") as f:
                data = f.read()
            file_name = file.split("/")[-1]
            if file_name == 'application.properties':
                a_result = {}
                # xx=xx
                pattern = r"(.)+=(.)+"
                data_list = data.strip().split("\n")
                for line in data_list:
                    if re.match(pattern, line):
                        line_list = line.split("=")
                        value = [p.strip() for p in line_list[1:]]
                        a_result[line_list[0].strip()] = "=".join(value)
                if a_result != {}:
                    result.append({file: [a_result]})
            elif file_name[-4:] == '.yml':
                a_result = {}
                # xx:xx
                pattern = r"(.)+:(.)+"
                data_list = data.strip().split("\n")
                for line in data_list:
                    if re.match(pattern, line):
                        line_list = line.split(":")
                        value = [p.strip() for p in line_list[1:]]
                        a_result[line_list[0].strip()] = ":".join(value)
                if a_result != {}:
                    result.append({file: [a_result]})
            elif file_name[-4:] == '.xml':
                a_result = {}
                # <property name="xxx"xxx
                pattern = r"<property name=\"(.)+\"(.)+"
                data_list = data.strip().split("\n")
                for line in data_list:
                    line = line.strip()
                    if re.match(pattern, line):
                        line_list = line.split("\"")
                        a_result[line_list[1].strip()] = line_list[3].strip()
                if a_result != {}:
                    result.append({file: [a_result]})
            elif file_name == '.bash_history':
                a_result = preict_command(file)
                result.append({file: a_result})
            elif file_name == 'authorized_keys' or file_name == 'id_rsa.pub':
                a_result = {}
                data_list = data.strip().split("\n")
                for line in data_list:
                    line = line.strip()
                    if line != "":
                        line_list = line.split(" ")
                        a_result[line_list[0].strip()] = line_list[1].strip()
                if a_result != {}:
                    result.append({file: [a_result]})
            elif file_name == 'id_rsa':
                a_result = {}
                # <property name="xxx"xxx
                pattern = r'-+BEGIN OPENSSH PRIVATE KEY-+(.)+-+END OPENSSH PRIVATE KEY-+'
                if re.match(pattern, data, re.S):
                    a_result["private_key"] = "\n".join(data.split("\n")[1:-2])
                if a_result != {}:
                    result.append({file: [a_result]})
            elif file_name == 'shadow':
                a_result = []
                data_list = data.strip().split("\n")
                for line in data_list:
                    dict = {}
                    line_list = line.split(":")
                    dict["username"] = line_list[0]
                    dict["password"] = line_list[1]
                    dict["modification_time"] = line_list[2]
                    dict["Minimum_validity_days"] = line_list[3]
                    dict["Maximum_validity_days"] = line_list[4]
                    dict["Grace_time_after_expiration"] = line_list[5]
                    dict["Account_Expiration_Date"] = line_list[6]
                    a_result.append(dict)
                if a_result != []:
                    result.append({file: a_result})
            elif file_name == 'passwd':
                a_result = []
                data_list = data.strip().split("\n")
                for line in data_list:
                    dict = {}
                    line_list = line.split(":")
                    dict["username"] = line_list[0]
                    dict["Password_Placeholder"] = line_list[1]
                    dict["UID"] = line_list[2]
                    dict["GID"] = line_list[3]
                    dict["Additional_basic_information"] = line_list[4]
                    dict["Home_Directory_Location"] = line_list[5]
                    dict["shell"] = line_list[6]
                    a_result.append(dict)
                if a_result != []:
                    result.append({file: a_result})
            elif file_name == 'token':
                a_result = {"token": data}
                result.append({file: [a_result]})
        # windows文件夹下注册表文件抽取
        elif ("windows" in file.split("/") or "windwos" in file.split("/")) and ("sam" in file or "system" in file):
            if file[-4:] == ".zip":
                temp_folder="temp_zip"
                if not os.path.exists(temp_folder):
                    os.makedirs(temp_folder)
                # 解压
                with zipfile.ZipFile(file, 'r') as zip_ref:
                    zip_ref.extractall(temp_folder)
                # 获取解压后文件
                z_file_ls=[]
                for z_root, z_dirs, z_files in os.walk(temp_folder):
                    z_root_file_ls = [os.path.join(z_root, z_file) for z_file in z_files]
                    z_file_ls += z_root_file_ls
                # 遍历，抽取
                hiv_dict_output = []
                for z_file in z_file_ls:
                    _, a_result = parse_and_save_hiv(z_file)
                    hiv_dict_output+=a_result
                shutil.rmtree(temp_folder)
            else:
                _, hiv_dict_output = parse_and_save_hiv(file)
            result.append({file:hiv_dict_output})
        # .pdf文本处理
        elif file[-4:]==".pdf":
            pdf_text=get_pdf_text(file)
            text_result.append({file:pdf_text})
        # .pptx文本处理
        elif file[-5:] == ".pptx":
            pptx_result,pptx_text = extract_text_from_pptx(file)
            text_result.append({file:pptx_text})
            if pptx_result!=[]:
                result.append({file: pptx_result})
        # .docx表格抽取，文本处理
        elif file[-5:] == ".docx":
            docx_result,docx_text = extract_text_from_docx(file)
            text_result.append({file:docx_text})
            if docx_result!=[]:
                result.append({file: docx_result})
        else:
            flag=0
        if flag==1:
            print("Extracting from %s" % file)

    with open(extract_path, "w", encoding='utf-8') as jsonl_file:
        for item in result:
            json_line = json.dumps(item, ensure_ascii=False, cls=DateEncoder)
            jsonl_file.write(json_line + "\n")
    with open(text_path, "w", encoding='utf-8') as jsonl_file:
        for item in text_result:
            json_line = json.dumps(item, ensure_ascii=False, cls=DateEncoder)
            jsonl_file.write(json_line + "\n")


class DateEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime.datetime):
            return obj.strftime("%Y-%m-%dT%H:%M:%S")
        else:
            return json.JSONEncoder.default(self, obj)



'''合并所有抽取结果'''
def conbine_all():
    pass


if __name__ == "__main__":
    # # 获取access token
    # access_token=get_access_token()
    # print(access_token)

    # 富文本路径
    path = "./data"
    # path = "./data/just_test"
    predict(path)

    # path="pptx_docx_text.jsonl"
    # data=[]
    # with open(path,"r") as f:
    #     for line in f:
    #         data.append(json.loads(line))
    # print("hh")